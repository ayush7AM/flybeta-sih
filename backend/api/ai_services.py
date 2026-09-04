"""
Real AI service for the Project Architect and Code Drishti features.
Uses the official google-genai SDK routed through Route429 proxy
for automatic API key rotation on rate limits.
"""
import os
import json
from google import genai
from google.genai import types
from pydantic import BaseModel, Field
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound

# ── Route429 Proxy Configuration ─────────────────────────────────────────
# Route429 sits between us and Gemini, managing a pool of API keys.
# It automatically rotates to the next key on HTTP 429 (rate limit).
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', 'route429-managed')
ROUTE429_BASE_URL = os.environ.get(
    'ROUTE429_BASE_URL',
    'https://route429.parth-ie-kalash.workers.dev/p/flybeta-sih'
)
ROUTE429_PROXY_SECRET = os.environ.get('ROUTE429_PROXY_SECRET', '')

# Initialize the Gemini client through Route429 proxy
# - api_key is a dummy value; Route429 injects the real key server-side
# - http_options overrides the base URL to point at the proxy
_http_options = {'base_url': ROUTE429_BASE_URL}
if ROUTE429_PROXY_SECRET:
    _http_options['headers'] = {'X-Proxy-Secret': ROUTE429_PROXY_SECRET}

try:
    client = genai.Client(
        api_key=GEMINI_API_KEY,
        http_options=_http_options,
    )
except Exception as e:
    print(f"[Route429] Failed to initialize Gemini client: {e}")
    client = None


# ── Pydantic Schemas for Structured Output ───────────────────────────────

class BlueprintStep(BaseModel):
    step_number: int
    tag: str = Field(description="Short uppercase category tag, e.g., DATA_PREP")
    title: str
    tasks: list[str]
    estimate: str

class BlueprintResponse(BaseModel):
    steps: list[BlueprintStep]

class ReviewFinding(BaseModel):
    id: int
    severity: str = Field(description="Must be one of: CRITICAL, WARNING, INFO, STYLE")
    category: str = Field(description="Must be one of: SECURITY, PERFORMANCE, LOGIC, STYLE")
    title: str
    description: str
    suggestion: str
    line_range: str = Field(description="Affected line range, e.g., '12-15' or '42'")

class ReviewResponse(BaseModel):
    findings: list[ReviewFinding]

class QuizQuestion(BaseModel):
    question: str
    options: list[str] = Field(description="Must contain exactly 4 options.")
    correct_index: int = Field(description="Must be 0, 1, 2, or 3 representing the index of the correct option.")

class QuizResponse(BaseModel):
    questions: list[QuizQuestion]


# ── AI Service Functions ──────────────────────────────────────────────────

def generate_project_blueprint(prompt):
    """
    Generate a step-by-step project blueprint from a user prompt using Gemini.
    """
    if not client:
        raise ValueError("Missing API Key. The Project Architect cannot run.")

    try:
        response = client.models.generate_content(
            model='gemini-3.6-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=(
                    "You are an elite software architect. Break down the user's project "
                    "idea into a logical, step-by-step development blueprint."
                ),
                response_mime_type="application/json",
                response_schema=BlueprintResponse,
                temperature=0.4,
            ),
        )
        
        # The response text is a JSON string matching BlueprintResponse
        data = json.loads(response.text)
        return data.get("steps", [])
        
    except Exception as e:
        print(f"Gemini API Error: {e}")
        raise ValueError(f"AI Generation failed: {str(e)}")


def generate_code_review(code, language="python"):
    """
    Generate a structured code review using Gemini.
    """
    if not client:
        raise ValueError("Missing API Key. The Code Reviewer cannot run.")

    try:
        prompt = f"Review the following {language} code:\n\n{code}"
        
        response = client.models.generate_content(
            model='gemini-3.6-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=(
                    "You are a strict, senior code reviewer. Analyze the provided code "
                    "for bugs, security flaws, performance issues, and style guide violations."
                ),
                response_mime_type="application/json",
                response_schema=ReviewResponse,
                temperature=0.2,
            ),
        )
        
        data = json.loads(response.text)
        return data.get("findings", [])
        
    except Exception as e:
        print(f"Gemini API Error: {e}")
        raise ValueError(f"AI Generation failed: {str(e)}")

def generate_video_quiz(video_id):
    """
    Extracts transcript from a YouTube video and generates a 3-question 
    multiple choice quiz using Gemini Structured Outputs.
    """
    if not client:
        raise ValueError("GEMINI_API_KEY is missing. Quiz generation cannot proceed.")

    # 1. Fetch Transcript
    try:
        transcript_list = YouTubeTranscriptApi().fetch(video_id)
        # Concatenate transcript texts
        raw_text = " ".join([t.text for t in transcript_list])
    except TranscriptsDisabled:
        raise ValueError("Captions are disabled for this video.")
    except NoTranscriptFound:
        raise ValueError("No transcript found for this video.")
    except Exception as e:
        raise ValueError(f"Could not extract transcript: {str(e)}")

    # 2. Prompt Gemini
    prompt = f"Create a 3-question multiple-choice quiz based on this transcript:\n\n{raw_text}"
    
    try:
        response = client.models.generate_content(
            model='gemini-3.6-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=(
                    "You are an expert educator. Based on the provided video transcript, "
                    "create a strict 3-question multiple-choice quiz. "
                    "Each question must have exactly 4 options and 1 correct_index."
                ),
                response_mime_type="application/json",
                response_schema=QuizResponse,
                temperature=0.3,
            ),
        )
        
        data = json.loads(response.text)
        return data.get("questions", [])
        
    except Exception as e:
        print(f"Gemini API Error: {e}")
        raise ValueError(f"AI Generation failed: {str(e)}")


def ask_oracle(message: str, history: list = None) -> str:
    """
    Acts as 'The Oracle', providing context-aware AI mentorship.
    """
    if not client:
        return "The Oracle is currently disconnected (Missing API Key)."

    history = history or []
    
    # Format the conversational history
    formatted_contents = []
    for msg in history:
        # Convert roles from our frontend format to Gemini format
        role = 'model' if msg.get('role') == 'ai' else 'user'
        formatted_contents.append({
            "role": role,
            "parts": [{"text": msg.get('content', '')}]
        })
    
    # Append the new user message
    formatted_contents.append({
        "role": "user",
        "parts": [{"text": message}]
    })
    
    system_instruction = (
        "You are The Oracle, a friendly, encouraging AI study buddy. "
        "Keep your response UNDER 40 WORDS. Be concise. "
        "ONLY output plain simple text with emojis. ABSOLUTELY NO MARKDOWN formatting. "
        "Do NOT use bold (**text**), italics, headers, or bullet points. "
        "Use emojis frequently to make the text colorful and engaging."
    )

    try:
        response = client.models.generate_content(
            model='gemini-3.6-flash',
            contents=formatted_contents,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.7,
            ),
        )
        return response.text
    except Exception as e:
        print(f"Gemini API Error (Oracle): {e}")
        raise ValueError(f"Oracle Generation failed: {str(e)}")


# ── Document-to-Quiz Pydantic Schemas ─────────────────────────────────

class DocQuizQuestion(BaseModel):
    question_number: int
    question_text: str
    options: list[str] = Field(min_length=4, max_length=4, description="Exactly 4 answer options")
    correct_answer: str = Field(description="The full text of the correct option")
    explanation: str = Field(description="Why this is correct, citing the source text")
    frac_quadrant: str = Field(description="One of: comp_statistical, comp_technical, comp_digital_governance, comp_behavioural")
    difficulty: str = Field(description="easy, intermediate, or advanced")

class DocQuizResponse(BaseModel):
    questions: list[DocQuizQuestion]


def extract_text_from_file(file_obj, filename):
    """
    Extract text content from a PDF or PPTX file.
    Returns the extracted text string.
    """
    ext = filename.lower().rsplit('.', 1)[-1] if '.' in filename else ''

    if ext == 'pdf':
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_obj) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return '\n\n'.join(text_parts)

    elif ext == 'pptx':
        from pptx import Presentation
        prs = Presentation(file_obj)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
        return '\n\n'.join(text_parts)

    else:
        raise ValueError(f"Unsupported file type: .{ext}. Only PDF and PPTX are supported.")


def generate_quiz_from_document(text, num_questions=5, difficulty='intermediate'):
    """
    Generate FRAC-tagged MCQs from extracted document text using Gemini.
    
    Args:
        text: The raw text extracted from the document.
        num_questions: Number of questions to generate (3-10).
        difficulty: One of 'easy', 'intermediate', 'advanced'.
    
    Returns:
        List of question dicts matching DocQuizQuestion schema.
    """
    if not client:
        raise ValueError("Missing API Key. Document quiz generation cannot run.")

    # Truncate to avoid context overflow
    max_chars = 8000
    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n[... document truncated for processing ...]"

    num_questions = max(3, min(10, int(num_questions)))
    difficulty = difficulty if difficulty in ('easy', 'intermediate', 'advanced') else 'intermediate'

    prompt = (
        f"Generate exactly {num_questions} multiple-choice questions at {difficulty} difficulty "
        f"based ONLY on the following document text. Do not use any external knowledge.\n\n"
        f"--- DOCUMENT TEXT ---\n{text}\n--- END DOCUMENT TEXT ---"
    )

    system_instruction = (
        "You are an expert assessment designer for India's Ministry of Statistics (MoSPI). "
        "Generate MCQs strictly from the provided document text. "
        "Each question must have exactly 4 options with 1 correct answer. "
        "The explanation MUST cite or reference the specific part of the document that supports the answer. "
        "Auto-tag each question to the most relevant MoSPI FRAC competency quadrant: "
        "comp_statistical (survey design, sampling, national accounts, SDG indicators), "
        "comp_technical (Python, R, SQL, GIS, AI/ML, data pipelines), "
        "comp_digital_governance (cybersecurity, data privacy, gov-cloud, DPI systems), "
        "comp_behavioural (leadership, communication, project management, ethics). "
        "If the content doesn't clearly fit any quadrant, default to comp_statistical."
    )

    try:
        response = client.models.generate_content(
            model='gemini-3.6-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                response_mime_type="application/json",
                response_schema=DocQuizResponse,
                temperature=0.3,
            ),
        )

        data = json.loads(response.text)
        return data.get("questions", [])

    except Exception as e:
        print(f"Gemini API Error (Doc Quiz): {e}")
        raise ValueError(f"Document quiz generation failed: {str(e)}")
